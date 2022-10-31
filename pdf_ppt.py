from fpdf import FPDF
import os
from PIL import Image
import fitz
import io

from tqdm import trange
from pptx import Presentation
from pptx.util import Cm
from pathlib import Path

def convert_to_pdf(path,dest,name):
    final=[]
    files = sorted(os.listdir(path))
    images=[]
    for i in files:
        if i.endswith(".png"):
            images.append(i)
    for i in images[1:]:
        img = Image.open(f"{path}/{i}")
        imgc = img.convert('RGB')
        final.append(imgc)
    img0=Image.open(f"{path}/{images[0]}")
    imgc0=img0.convert("RGB")
    imgc0.save(f"{dest}/{name}.pdf",save_all=True,append_images=final)

def convert_to_ppt(input_path, output_file, name,resolution=300, start_page=0, page_count=None,quiet=True):
    convert_to_pdf(input_path,os.getcwd(),"temp")
    pdf_file = "temp.pdf"
    doc = fitz.open("temp.pdf")
    if not quiet:
        print(pdf_file, 'contains', doc.page_count, 'slides')

    if page_count is None:
        page_count = doc.page_count

    # transformation matrix: slide to pixmap
    zoom = resolution / 72
    matrix = fitz.Matrix(zoom, zoom, 0)

    # create pptx presentation
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    # configure presentation aspect ratio
    page = doc.load_page(0)
    aspect_ratio = page.rect.width / page.rect.height
    prs.slide_width = int(prs.slide_height * aspect_ratio)
    
    # create page iterator
    if not quiet:
        page_iter = trange(start_page, start_page + page_count)
    else:
        print(start_page,  page_count)
        page_iter = range(start_page, start_page + page_count)

    # iterate over slides
    for page_no in page_iter:
        page = doc.load_page(page_no)

        # write slide as a pixmap
        pixmap = page.get_pixmap(matrix=matrix)
        image_data = pixmap.tobytes(output='PNG')
        image_file = io.BytesIO(image_data)
    
        # add a slide
        slide = prs.slides.add_slide(blank_slide_layout)    
        left = top = Cm(0)
        slide.shapes.add_picture(image_file, left, top, height=prs.slide_height)

    if output_file is None:
        if name is None:
            output_file = Path(pdf_file).with_suffix('.pptx')
        else:
            output_file = Path(pdf_file).with_suffix(f'{name}.pptx')
    else:
        output_file = output_file+"\\"+name+".pptx"
    # save presentation
    print(output_file)
    prs.save(output_file)
    doc.close()
    os.remove(os.getcwd()+"\\temp.pdf")

# convert_to_ppt("test_photos/test6","test.pptx")

def carousel_to_pdf(pages,path,dest,name):
    final=[]
    images = sorted(os.listdir(path))
    ind=1
    for i in images[1:]:
        if ind+1 in pages[1:]:
            img = Image.open(f"{path}/{i}")
            imgc = img.convert('RGB')
            final.append(imgc)
        ind+=1
    img0=Image.open(f"{path}/{images[pages[0]]}")
    imgc0=img0.convert("RGB")
    imgc0.save(f"{dest}/{name}.pdf",save_all=True,append_images=final)


# carousel_to_pdf([1,2,3,4,5],"test_photos/test6",os.getcwd(),"test")
# input_path="test_photos/test6"
# convert_to_pdf(input_path,os.getcwd(),"temp")

# convert_to_pdf("test_photos/test6","test_pdf","test6")